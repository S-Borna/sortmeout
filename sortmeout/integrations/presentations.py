"""
Presentation builder — generate PowerPoint/Keynote presentations.

Capabilities:
    - Create presentations from AI-generated outlines
    - Add slides with title, content, bullet points
    - Support for different slide layouts
    - Export to PowerPoint (.pptx) format
    - Open in Keynote after creation

Uses python-pptx for generation (no Keynote AppleScript needed).
Falls back to Keynote AppleScript if python-pptx not available.
"""

from __future__ import annotations

import os
import subprocess
from pathlib import Path
from datetime import datetime
from typing import Optional, Dict, List, Any

from sortmeout.utils.logger import get_logger

logger = get_logger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class PresentationBuilder:
    """Build presentations from structured data."""

    def __init__(self, output_dir: Optional[str] = None):
        """Initialize builder.

        Args:
            output_dir: Where to save presentations. Defaults to ~/Documents/Presentations/
        """
        self.output_dir = output_dir or os.path.expanduser("~/Documents/Presentations")
        os.makedirs(self.output_dir, exist_ok=True)

    def create_presentation(
        self,
        title: str,
        slides: List[Dict[str, Any]],
        subtitle: Optional[str] = None,
        author: Optional[str] = None,
        open_after: bool = True,
    ) -> Dict[str, Any]:
        """Create a presentation from structured slide data.

        Args:
            title: Presentation title (also used for filename)
            slides: List of slide dicts, each with:
                - title: str — slide title
                - content: str or list — body text or bullet points
                - layout: str — "title", "bullets", "two_column", "blank"
                - notes: str — speaker notes (optional)
            subtitle: Subtitle for title slide
            author: Author name
            open_after: Open in Keynote after saving

        Returns:
            Dict with file path and success status.
        """
        if HAS_PPTX:
            return self._create_with_pptx(title, slides, subtitle, author, open_after)
        else:
            return self._create_with_keynote(title, slides, subtitle, open_after)

    def _create_with_pptx(
        self,
        title: str,
        slides: List[Dict],
        subtitle: Optional[str],
        author: Optional[str],
        open_after: bool,
    ) -> Dict[str, Any]:
        """Create presentation using python-pptx."""
        prs = Presentation()

        # Set slide dimensions (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title
        if subtitle and slide.placeholders[1]:
            slide.placeholders[1].text = subtitle

        # Content slides
        for slide_data in slides:
            slide_title = slide_data.get("title", "")
            content = slide_data.get("content", "")
            layout_type = slide_data.get("layout", "bullets")
            notes = slide_data.get("notes", "")

            if layout_type == "title":
                # Section header
                layout = prs.slide_layouts[2]  # Section Header
                s = prs.slides.add_slide(layout)
                s.shapes.title.text = slide_title
                if content and len(s.placeholders) > 1:
                    s.placeholders[1].text = (
                        content if isinstance(content, str) else "\n".join(content)
                    )

            elif layout_type == "bullets":
                layout = prs.slide_layouts[1]  # Title and Content
                s = prs.slides.add_slide(layout)
                s.shapes.title.text = slide_title

                body = s.placeholders[1]
                tf = body.text_frame
                tf.clear()

                if isinstance(content, list):
                    for i, bullet in enumerate(content):
                        if i == 0:
                            tf.text = bullet
                        else:
                            p = tf.add_paragraph()
                            p.text = bullet
                            p.level = 0
                elif isinstance(content, str):
                    tf.text = content

            elif layout_type == "two_column":
                # Use blank + manual text boxes for two columns
                layout = prs.slide_layouts[5]  # Blank
                s = prs.slides.add_slide(layout)

                # Title
                from pptx.util import Inches, Pt

                txBox = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
                tf = txBox.text_frame
                tf.text = slide_title
                tf.paragraphs[0].font.size = Pt(28)
                tf.paragraphs[0].font.bold = True

                # Two columns
                if isinstance(content, list) and len(content) >= 2:
                    # Left column
                    left = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5))
                    left.text_frame.text = (
                        content[0] if isinstance(content[0], str) else "\n".join(content[0])
                    )

                    # Right column
                    right = s.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(5))
                    right.text_frame.text = (
                        content[1] if isinstance(content[1], str) else "\n".join(content[1])
                    )

            else:
                # Blank with just text
                layout = prs.slide_layouts[5]
                s = prs.slides.add_slide(layout)
                if content:
                    txBox = s.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(5.5))
                    tf = txBox.text_frame
                    tf.text = content if isinstance(content, str) else "\n".join(content)

            # Speaker notes
            if notes:
                s.notes_slide.notes_text_frame.text = notes

        # Save
        safe_title = "".join(c for c in title if c.isalnum() or c in " -_").strip()
        filename = f"{safe_title}.pptx"
        filepath = os.path.join(self.output_dir, filename)

        # Handle existing file
        counter = 1
        while os.path.exists(filepath):
            filepath = os.path.join(self.output_dir, f"{safe_title} ({counter}).pptx")
            counter += 1

        prs.save(filepath)

        # Open in Keynote
        if open_after:
            subprocess.Popen(["open", filepath])

        return {
            "success": True,
            "path": filepath,
            "slides_count": len(slides) + 1,  # +1 for title slide
            "format": "pptx",
        }

    def _create_with_keynote(
        self,
        title: str,
        slides: List[Dict],
        subtitle: Optional[str],
        open_after: bool,
    ) -> Dict[str, Any]:
        """Create presentation using Keynote AppleScript (fallback)."""
        # Build slide content as AppleScript
        slide_scripts = []
        for i, slide_data in enumerate(slides):
            slide_title = slide_data.get("title", "").replace('"', '\\"')
            content = slide_data.get("content", "")
            if isinstance(content, list):
                body = "\\n".join(c.replace('"', '\\"') for c in content)
            else:
                body = content.replace('"', '\\"').replace("\n", "\\n")

            slide_scripts.append(
                f"""
            tell slide {i + 2}
                set object text of default title item to "{slide_title}"
                set object text of default body item to "{body}"
            end tell
            """
            )

        slides_code = "\n".join(slide_scripts)

        subtitle_line = ""
        if subtitle:
            escaped_subtitle = subtitle.replace('"', '\\"')
            subtitle_line = f'set object text of default body item to "{escaped_subtitle}"'

        script = f"""
        tell application "Keynote"
            activate
            set newDoc to make new document with properties {{document theme:theme "Gradient"}}
            tell newDoc
                tell slide 1
                    set object text of default title item to "{title}"
                    {subtitle_line}
                end tell

                -- Add content slides
                repeat {len(slides)} times
                    make new slide
                end repeat

                {slides_code}
            end tell
        end tell
        """

        try:
            subprocess.run(["osascript", "-e", script], capture_output=True, text=True, timeout=30)
            return {
                "success": True,
                "path": "(Keynote document - save manually)",
                "slides_count": len(slides) + 1,
                "format": "keynote",
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def create_from_outline(self, outline: str, title: Optional[str] = None) -> Dict[str, Any]:
        """Create a presentation from a text outline.

        The outline format:
            # Presentation Title
            ## Slide Title
            - Bullet point 1
            - Bullet point 2
            ## Another Slide
            - Point A
            - Point B

        This is what the AI will generate and pass here.
        """
        lines = outline.strip().split("\n")
        prs_title = title or "Presentation"
        slides = []
        current_slide = None

        for line in lines:
            line = line.strip()
            if not line:
                continue

            if line.startswith("# ") and not line.startswith("## "):
                prs_title = line[2:].strip()
            elif line.startswith("## "):
                if current_slide:
                    slides.append(current_slide)
                current_slide = {
                    "title": line[3:].strip(),
                    "content": [],
                    "layout": "bullets",
                }
            elif line.startswith("- ") or line.startswith("* "):
                if current_slide:
                    current_slide["content"].append(line[2:].strip())
            elif current_slide:
                # Plain text — append to content
                if current_slide["content"]:
                    current_slide["content"][-1] += " " + line
                else:
                    current_slide["content"].append(line)

        if current_slide:
            slides.append(current_slide)

        return self.create_presentation(prs_title, slides)
